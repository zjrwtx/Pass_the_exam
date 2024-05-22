# app/main.py
from fastapi import FastAPI, File, UploadFile, Form
from typing import List
from PyPDF2 import PdfReader
from pptx import Presentation
import os
from io import BytesIO
from fastapi.middleware.cors import CORSMiddleware
from dotenv import load_dotenv
import requests
def maketopics(text, typeoftopic, numoftopic):
    # 加载环境变量文件
    load_dotenv()
    API_KEY = os.getenv('API_KEY')
    API_url = os.getenv('API_url')

    # 设置请求头
    headers = {
        'Content-Type': 'application/json',
        'Authorization': f"Bearer {API_KEY}",
    }

    # 设置要发送的数据
    json_data = {
        'model': 'yi-large',
        'messages': [
            {   
                'role': 'system',
                'content': '你是一个专业的出题者，能对用户输入的内容生成详细的高质量的题目，且带上答案与解释',
            },
            {
                'role': 'user',
                'content': f"请你把我输入的内容生成{numoftopic}道详细的高质量的{typeoftopic}类型题目，且带上答案与解释: {text}",
            },
        ],
        'temperature': 0.3,
         "max_tokens": 10000,
         "top_p": 1.0,
    }

    # 发送请求
    response = requests.post(API_url, headers=headers, json=json_data)

    # 打印响应内容
    try:
        output = response.json()['choices'][0]['message']['content']
        print(output)
        return output

    except (KeyError, IndexError, AttributeError):
        print("无法获取响应内容或响应格式不正确。")
        return "出错啦！ 请联系和push一下开发者解决问题"




def mindmap(text):
    # 加载环境变量文件
    load_dotenv()
    API_KEY = os.getenv('API_KEY')
    API_url = os.getenv('API_url')

    # 设置请求头
    headers = {
        'Content-Type': 'application/json',
        'Authorization': f"Bearer {API_KEY}",
    }

    # 设置要发送的数据
    json_data = {
        'model': 'yi-large',
        'messages': [
            {   
                'role': 'system',
                'content': '你是一个专业的思维导图专家，能对用户输入的内容生成详细的高质量的markdown格式的思维导图',
            },
            {
                'role': 'user',
                'content': f"请你把我输入的内容生成详细的高质量的markdown格式的思维导图{text}",
            },
        ],
        'temperature': 0.3,
         "max_tokens": 10000,
         "top_p": 1.0,
    }

    # 发送请求
    response = requests.post(API_url, headers=headers, json=json_data)

    # 打印响应内容
    try:
        output = response.json()['choices'][0]['message']['content']
        print(output)
        return output

    except (KeyError, IndexError, AttributeError):
        print("无法获取响应内容或响应格式不正确。")
        return "出错啦！ 请联系和push一下开发者解决问题"

app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.post("/topics")
async def upload_files(
    files: List[UploadFile] = File(...),
    typeoftopic: str = Form(...),
    numoftopic: str = Form(...)
):
    responses = []
    for file in files:
        content = await file.read()
        filename = file.filename
        ext = os.path.splitext(filename)[1].lower()

        if ext == ".pdf":
            extracted_text = extract_text_from_pdf(content)
        elif ext == ".pptx" or ext == ".ppt":
            extracted_text = extract_text_from_pptx(content)
        else:
            extracted_text = "Unsupported file type"

        api_response = maketopics(extracted_text, typeoftopic, numoftopic)
        responses.append({"filename": filename, "api_response": api_response})

    return {"responses": responses}
@app.post("/mindmap")
async def upload_files(
    files: List[UploadFile] = File(...),

):
    responses = []
    for file in files:
        content = await file.read()
        filename = file.filename
        ext = os.path.splitext(filename)[1].lower()

        if ext == ".pdf":
            extracted_text = extract_text_from_pdf(content)
        elif ext == ".pptx" or ext == ".ppt":
            extracted_text = extract_text_from_pptx(content)
        else:
            extracted_text = "Unsupported file type"

        api_response = mindmap(extracted_text)
        responses.append({"filename": filename, "api_response": api_response})

    return {"responses": responses}


def extract_text_from_pdf(content):
    pdf_file = BytesIO(content)
    pdf_reader = PdfReader(pdf_file)
    text = ''
    for page_num in range(len(pdf_reader.pages)):
        page = pdf_reader.pages[page_num]
        text += page.extract_text()
    return text

def extract_text_from_pptx(content):
    pptx_file = BytesIO(content)
    prs = Presentation(pptx_file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return "\n".join(text_runs)

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
