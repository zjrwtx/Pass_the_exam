from typing import Optional
from sqlalchemy.orm import Session
from passlib.context import CryptContext
from jose import JWTError, jwt
from datetime import datetime, timedelta, timezone
from app.models import UserModel, Feedback
from app.schemas import TokenData, User, FeedbackCreate
from fastapi import Depends, HTTPException, status
from fastapi.security import OAuth2PasswordBearer
import requests
from sqlalchemy import create_engine
import os
from dotenv import load_dotenv
from sqlalchemy.orm import sessionmaker

load_dotenv()
ALGORITHM = "HS256"
DATABASE_URL = os.getenv("DATABASE_URL")
SECRET_KEY = os.getenv("SECRET_KEY")
API_KEY = os.getenv('API_KEY')
API_url = os.getenv('API_url')

engine = create_engine(DATABASE_URL)
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)

pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token")

token_blacklist = []

def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()
def create_feedback(feedback: FeedbackCreate):
    db = SessionLocal()
    db_feedback = Feedback(
        name=feedback.name,
        email=feedback.email,
        feedback=feedback.feedback,
        created_at=datetime.utcnow()
    )
    db.add(db_feedback)
    db.commit()
    db.refresh(db_feedback)
    return db_feedback

def read_feedback():
    db = SessionLocal()
    feedback_list = db.query(Feedback).all()
    return feedback_list

def verify_password(plain_password, hashed_password):
    return pwd_context.verify(plain_password, hashed_password)

def get_password_hash(password):
    return pwd_context.hash(password)

def get_user(db: Session, username: str):
    return db.query(UserModel).filter(UserModel.username == username).first()

def authenticate_user(db: Session, username: str, password: str):
    user = get_user(db, username)
    if not user or not verify_password(password, user.hashed_password):
        return False
    return user

def create_access_token(data: dict, expires_delta: Optional[timedelta] = None):
    to_encode = data.copy()
    expire = datetime.now(timezone.utc) + expires_delta if expires_delta else datetime.now(timezone.utc) + timedelta(minutes=15)
    to_encode.update({"exp": expire})
    encoded_jwt = jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)
    return encoded_jwt

async def get_current_user(token: str = Depends(oauth2_scheme), db: Session = Depends(get_db)):
    if token in token_blacklist:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Token has been revoked", headers={"WWW-Authenticate": "Bearer"})
    
    credentials_exception = HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Could not validate credentials", headers={"WWW-Authenticate": "Bearer"})
    
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        username: str = payload.get("sub")
        if username is None:
            raise credentials_exception
        token_data = TokenData(username=username)
    except JWTError:
        raise credentials_exception
    
    user = get_user(db, username=token_data.username)
    if user is None:
        raise credentials_exception
    return user

async def get_current_active_user(current_user: User = Depends(get_current_user)):
    if current_user.disabled:
        raise HTTPException(status_code=400, detail="Inactive user")
    return current_user

def maketopics(text, typeoftopic, numoftopic):
    headers = {'Content-Type': 'application/json', 'Authorization': f"Bearer {API_KEY}"}
    json_data = {
        'model': 'step-1-128k',
        'messages': [
            {'role': 'system', 'content': '你是一个专业的出题者，能对用户输入的内容生成详细的高质量的题目，且带上答案与解释'},
            {'role': 'user', 'content': f"请你把我输入的内容生成{numoftopic}道详细的高质量的{typeoftopic}类型题目，且带上答案与解释: {text}"},
        ],
        'temperature': 0.3,
        "max_tokens": 10000,
        "top_p": 1.0,
    }
    response = requests.post(API_url, headers=headers, json=json_data)
    try:
        return response.json()['choices'][0]['message']['content']
    except (KeyError, IndexError, AttributeError) as e:
        return f"出错啦！ 请联系和push一下开发者解决问题。错误详情: {str(e)}"
    except requests.exceptions.RequestException as e:
        return f"请求错误！ 请检查网络连接或API服务状态。错误详情: {str(e)}"

def mindmap(text):
    headers = {'Content-Type': 'application/json', 'Authorization': f"Bearer {API_KEY}"}
    json_data = {
        'model': 'step-1-128k',
        'messages': [
            {'role': 'system', 'content': '你是一个专业的思维导图专家，能对用户输入的内容生成详细的高质量的markdown格式的思维导图'},
            {'role': 'user', 'content': f"请你把我输入的内容生成详细的高质量的markdown格式的思维导图{text}"},
        ],
        'temperature': 0.3,
        "max_tokens": 10000,
        "top_p": 1.0,
    }
    response = requests.post(API_url, headers=headers, json=json_data)
    try:
        return response.json()['choices'][0]['message']['content']
    except (KeyError, IndexError, AttributeError):
        return "出错啦！ 请联系和push一下开发者解决问题"

def examkeypoints(text):
    headers = {'Content-Type': 'application/json', 'Authorization': f"Bearer {API_KEY}"}
    json_data = {
        'model': 'step-1-128k',
        'messages': [
            {'role': 'system', 'content': '你是一个考试重点大纲助手，请你仔细浏览文档内容后，抽取里面的可能得重点考点出来，且按照重要程度去标星，最后列成详细的考点大纲'},
            {'role': 'user', 'content': f"请你仔细浏览文档内容后，抽取里面的可能得重点考点出来，且按照重要程度去标星，最后列成详细的考点大纲：{text}"},
        ],
        'temperature': 0.3,
        "max_tokens": 10000,
        "top_p": 1.0,
    }
    response = requests.post(API_url, headers=headers, json=json_data)
    try:
        return response.json()['choices'][0]['message']['content']
    except (KeyError, IndexError, AttributeError):
        return "出错啦！ 请联系和push一下开发者解决问题"
    
def exampaper(text,typeoftopic,numoftopic):
    headers = {'Content-Type': 'application/json', 'Authorization': f"Bearer {API_KEY}"}
    json_data = {
        'model': 'step-1-128k',
        'messages': [
            {'role': 'system', 'content': f'你是一个期末考试试卷生成助手，请你仔细浏览文档内容后，抽取里面的可能得重点考点出来，且按照重要程度去生成{numoftopic}道不同类型的题目：主要包括{typeoftopic}等题目，最后生成一份详细且格式美观的期末考试试卷'},
            {'role': 'user', 'content': f"请你仔细浏览文档内容后，抽取里面的可能得重点考点出来，且按照重要程度去生成{numoftopic}道不同类型的题目：主要包括{typeoftopic}等题目，最后生成一份详细且格式美观的期末考试试卷：{text}"},
        ],
        'temperature': 0.3,
        "max_tokens": 10000,
        "top_p": 1.0,
    }
    response = requests.post(API_url, headers=headers, json=json_data)
    try:
        return response.json()['choices'][0]['message']['content']
    except (KeyError, IndexError, AttributeError):
        return "出错啦！ 请联系和push一下开发者解决问题"

def extract_text_from_pdf(content):
    from PyPDF2 import PdfReader
    from io import BytesIO
    
    pdf_file = BytesIO(content)
    pdf_reader = PdfReader(pdf_file)
    text = ''
    for page_num in range(len(pdf_reader.pages)):
        page = pdf_reader.pages[page_num]
        text += page.extract_text()
    return text

def extract_text_from_pptx(content):
    from pptx import Presentation
    from io import BytesIO
    
    pptx_file = BytesIO(content)
    prs = Presentation(pptx_file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return "\n".join(text_runs)

def clean_up(paths):
    import os
    import shutil

    for path in paths:
        if os.path.isdir(path):
            shutil.rmtree(path)
        elif os.path.isfile(path):
            os.remove(path)
